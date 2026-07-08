# -*- coding: utf-8 -*-
"""네트워크 구성도 PPTX 자동 생성 — 수집 데이터 → 고객 전달용 슬라이드.

렌더: Cairo/Graphviz 없이 Pillow로 다이어그램 PNG를 그리고 python-pptx로 조립.
좌표는 topology.build_topology()의 노드/링크를 계층/구역별로 자체 배치.
슬라이드: 표지 → 서버실 구성도 → TPS 구역 개요 → 구역별 상세(N장) → 변경이력.
"""
import io
import os
import re
from datetime import datetime

from . import db, topology

# SK AX 브랜드 컬러
SK_RED = (234, 0, 44)
SK_ORANGE = (255, 122, 0)
DARK = (15, 23, 42)
GRAY = (100, 116, 139)
OK = (34, 197, 94)
WARN = (245, 158, 11)
BAD = (239, 68, 68)

_KIND_COLOR = {
    "Firewall": SK_RED, "BackBone": (168, 85, 247), "L3 Switch": (139, 92, 246),
    "L4 Switch": SK_ORANGE, "L2 Switch": (20, 184, 166), "Server": (59, 130, 246),
}
_KIND_LABEL = {
    "Firewall": "FW", "BackBone": "Core", "L3 Switch": "L3",
    "L4 Switch": "L4", "L2 Switch": "L2", "Server": "SRV",
}


def _logo_path():
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    p = os.path.join(here, "web", "static", "brand", "skax_logo.png")
    return p if os.path.exists(p) else None


def _font(size, bold=False):
    from PIL import ImageFont
    # Windows 기본 폰트 시도 → 실패 시 기본
    for name in (("malgunbd.ttf" if bold else "malgun.ttf"),
                 ("arialbd.ttf" if bold else "arial.ttf")):
        try:
            return ImageFont.truetype(name, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _status_color(n):
    if n.get("alert") == "critical" or n.get("status") == "failed" or n.get("reachable") is False:
        return BAD
    if n.get("alert") == "warning":
        return WARN
    if n.get("status") == "done":
        return OK
    return GRAY


def _node_dtype(n):
    return n.get("device_type") or topology.infer_role(n.get("name"), "") or "L2 Switch"


def _draw_device(draw, x, y, w, h, node):
    """장비 아이콘(사각형 실루엣 + 종류색 띠 + 상태색 테두리) + 라벨."""
    dt = _node_dtype(node)
    kc = _KIND_COLOR.get(dt, GRAY)
    sc = _status_color(node)
    draw.rounded_rectangle([x, y, x + w, y + h], radius=8, fill=(30, 41, 59), outline=sc, width=3)
    draw.rounded_rectangle([x, y, x + 8, y + h], radius=3, fill=kc)
    f = _font(15, bold=True); fs = _font(12)
    label = _KIND_LABEL.get(dt, "")
    draw.text((x + 16, y + 8), (label + "  " if label else "") + (node.get("name") or "")[:22],
              fill=(241, 245, 249), font=f)
    draw.text((x + 16, y + h - 22), (node.get("ip") or ""), fill=(148, 163, 184), font=fs)
    # 상태 점
    draw.ellipse([x + w - 20, y + 8, x + w - 8, y + 20], fill=sc)


def _curve(draw, a, b, color, width, dashed=False):
    """두 점 사이 부드러운 세로 곡선(간이 베지어)."""
    x1, y1 = a; x2, y2 = b
    steps = 24
    prev = None
    for i in range(steps + 1):
        t = i / steps
        mt = 1 - t
        cy1 = (y1 + y2) / 2
        # quadratic-ish through midpoints
        xx = mt * mt * x1 + 2 * mt * t * ((x1 + x2) / 2) + t * t * x2
        yy = mt * mt * y1 + 2 * mt * t * cy1 + t * t * y2
        if prev and not (dashed and i % 2 == 0):
            draw.line([prev, (xx, yy)], fill=color, width=width)
        prev = (xx, yy)


def _render_layered_png(nodes, links, title, width=1280, height=720):
    """계층형(rank별) 배치 다이어그램 PNG 생성. rank_fn으로 계층 결정."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (width, height), (11, 18, 36))
    draw = ImageDraw.Draw(img)
    # 제목
    draw.text((28, 18), title, fill=(226, 232, 240), font=_font(22, bold=True))

    def rank(n):
        dt = _node_dtype(n)
        if dt == "Firewall" or n.get("kind") == "fw":
            return 0
        if dt == "BackBone":
            return 1
        if dt in ("L3 Switch", "L4 Switch"):
            return 2
        return 3
    by_rank = {}
    for n in nodes:
        by_rank.setdefault(rank(n), []).append(n)
    ranks = sorted(by_rank)
    NW, NH = 210, 54
    top_margin = 70
    lane_h = (height - top_margin - 30) / max(1, len(ranks))
    pos = {}
    for ri, r in enumerate(ranks):
        row = sorted(by_rank[r], key=lambda n: (n.get("name") or ""))
        gap = width / (len(row) + 1)
        y = top_margin + ri * lane_h + (lane_h - NH) / 2
        for ci, n in enumerate(row):
            pos[n["id"]] = (gap * (ci + 1) - NW / 2, y)
    # 링크
    byid = {n["id"]: n for n in nodes}
    for l in links:
        if l["a"] not in pos or l["b"] not in pos:
            continue
        pa = pos[l["a"]]; pb = pos[l["b"]]
        a = (pa[0] + NW / 2, pa[1] + NH); b = (pb[0] + NW / 2, pb[1])
        if pa[1] > pb[1]:
            a, b = (pb[0] + NW / 2, pb[1] + NH), (pa[0] + NW / 2, pa[1])
        cdp = l.get("source") == "cdp/lldp"
        _curve(draw, a, b, (56, 189, 248) if cdp else (148, 163, 184),
               3 if cdp else 2, dashed=not l.get("mutual"))
        # 포트 라벨
        pa_txt = (l.get("a_port") or "")[:12]
        if pa_txt:
            mx, my = (a[0] + b[0]) / 2, (a[1] + b[1]) / 2
            draw.text((mx - 20, my - 12), pa_txt, fill=(125, 211, 252), font=_font(11))
    # 노드
    for n in nodes:
        if n["id"] in pos:
            x, y = pos[n["id"]]
            _draw_device(draw, x, y, NW, NH, n)
    buf = io.BytesIO(); img.save(buf, "PNG"); buf.seek(0)
    return buf


def _zone_of(n):
    if n.get("depth") == 0:
        return "코어/백본"
    return (n.get("group") or "").strip() or "미지정"


def _is_core(n):
    dt = _node_dtype(n)
    return n.get("kind") == "fw" or dt in ("Firewall", "BackBone", "L3 Switch", "L4 Switch", "Server")


def _stats(nodes):
    ok = sum(1 for n in nodes if _status_color(n) == OK)
    warn = sum(1 for n in nodes if _status_color(n) == WARN)
    bad = sum(1 for n in nodes if _status_color(n) == BAD)
    return {"total": len(nodes), "ok": ok, "warn": warn, "bad": bad}


def build_pptx(db_path, customer="", generated=None):
    """구성도 PPTX 바이트 반환. customer=고객사명, generated=생성일자(YYYY-MM-DD)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    generated = generated or datetime.now().strftime("%Y-%m-%d")
    topo = topology.build_topology(db_path)
    nodes = topo.get("nodes", []); links = topo.get("links", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height
    logo = _logo_path()

    def _rgb(t): return RGBColor(*t)

    def _add_title_bar(slide, text):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.9))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(DARK); bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.6))
        p = tb.text_frame.paragraphs[0]; run = p.add_run(); run.text = text
        run.font.size = Pt(24); run.font.bold = True; run.font.color.rgb = _rgb((255, 255, 255))
        if logo:
            slide.shapes.add_picture(logo, SW - Inches(2.3), Inches(0.12), height=Inches(0.62))

    def _add_png(slide, buf, top=Inches(1.0)):
        slide.shapes.add_picture(buf, Inches(0.3), top, width=SW - Inches(0.6))

    def _stat_box(slide, st, top=Inches(6.7)):
        tb = slide.shapes.add_textbox(Inches(0.4), top, SW - Inches(0.8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        for label, val, col in [("전체 ", st["total"], DARK), ("  정상 ", st["ok"], OK),
                                 ("  경고 ", st["warn"], WARN), ("  장애 ", st["bad"], BAD)]:
            r = p.add_run(); r.text = "%s%d" % (label, val)
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = _rgb(col)

    # 1) 표지
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = _rgb((248, 250, 252)); bg.line.fill.background()
    if logo:
        s.shapes.add_picture(logo, Inches(4.8), Inches(1.4), width=Inches(3.7))
    tb = s.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "네트워크 구성도"; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = _rgb(DARK)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = (customer or "고객사") + "  ·  " + generated
    r2.font.size = Pt(20); r2.font.color.rgb = _rgb(SK_RED)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "SK AX USA  |  NetDash 자동 생성"
    r3.font.size = Pt(12); r3.font.color.rgb = _rgb(GRAY)

    # 2) 서버실 구성도(코어 장비)
    core = [n for n in nodes if _is_core(n)]
    core_ids = {n["id"] for n in core}
    core_links = [l for l in links if l["a"] in core_ids and l["b"] in core_ids]
    s = prs.slides.add_slide(blank); _add_title_bar(s, "서버실 구성도 (코어 인프라)")
    _add_png(s, _render_layered_png(core, core_links, "Firewall → Core → Distribution → Access"))
    _stat_box(s, _stats(core))

    # 3) TPS 구역 개요
    access = [n for n in nodes if not _is_core(n)]
    zones = {}
    for n in access:
        zones.setdefault(_zone_of(n), []).append(n)
    s = prs.slides.add_slide(blank); _add_title_bar(s, "TPS 구역 개요")
    zkeys = sorted(zones)
    cols = 4
    bw, bh = Inches(3.0), Inches(1.4)
    x0, y0 = Inches(0.4), Inches(1.2)
    for i, z in enumerate(zkeys):
        st = _stats(zones[z])
        col = i % cols; row = i // cols
        bx = Emu(int(x0) + col * int(Inches(3.2)))
        by = Emu(int(y0) + row * int(Inches(1.6)))
        box = s.shapes.add_shape(1, bx, by, bw, bh)
        worst = BAD if st["bad"] else (WARN if st["warn"] else OK)
        box.fill.solid(); box.fill.fore_color.rgb = _rgb((30, 41, 59))
        box.line.color.rgb = _rgb(worst); box.line.width = Pt(2.5)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = z[:24]
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = _rgb((241, 245, 249))
        p2 = tf.add_paragraph(); r2 = p2.add_run()
        r2.text = "장비 %d · 정상 %d · 경고 %d · 장애 %d" % (st["total"], st["ok"], st["warn"], st["bad"])
        r2.font.size = Pt(11); r2.font.color.rgb = _rgb(worst)
    if not zkeys:
        tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        tb.text_frame.paragraphs[0].add_run().text = "TPS(액세스) 장비가 없습니다."

    # 4) 구역별 상세(구역 하나당 슬라이드, 40대↑이면 분할)
    for z in zkeys:
        zn = zones[z]
        # 업링크(코어) 고스트 포함
        znids = {n["id"] for n in zn}
        ghost_ids = set()
        for l in links:
            if l["a"] in znids and l["b"] not in znids:
                ghost_ids.add(l["b"])
            if l["b"] in znids and l["a"] not in znids:
                ghost_ids.add(l["a"])
        ghosts = [n for n in nodes if n["id"] in ghost_ids]
        chunks = [zn[i:i + 40] for i in range(0, len(zn), 40)] or [[]]
        for ci, chunk in enumerate(chunks):
            show = ghosts + chunk
            show_ids = {n["id"] for n in show}
            zl = [l for l in links if l["a"] in show_ids and l["b"] in show_ids]
            title = "구역 상세 — %s%s" % (z, ("  (%d/%d)" % (ci + 1, len(chunks)) if len(chunks) > 1 else ""))
            s = prs.slides.add_slide(blank); _add_title_bar(s, title)
            _add_png(s, _render_layered_png(show, zl, z))
            _stat_box(s, _stats(chunk))

    out = io.BytesIO(); prs.save(out); out.seek(0)
    fname = "%s_네트워크구성도_%s.pptx" % (
        re.sub(r"[^\w가-힣.-]", "_", customer or "고객사"),
        generated.replace("-", ""))
    return out.getvalue(), fname
